import os
from datetime import datetime
from pptx import Presentation
from pptx.exc import PackageNotFoundError


def get_presentation(presentation_selection):
    """This function is the only place where the presentations actually get opened"""
    if presentation_selection:
        try:

            with open(presentation_selection, 'rb') as presentation_file:
                presentation_file = Presentation(presentation_selection)

                presentation_name: str = os.path.split(
                    presentation_selection)[1]
                presentation_path: str = os.path.split(
                    presentation_selection)[0]

                slide_number: int = len(presentation_file.slides)
                hidden_slides: int = get_hidden_slides(presentation_file)

                presentation_created: str = get_time(0, presentation_selection)
                presentation_modified: str = get_time(
                    1, presentation_selection)

                all_text: list = get_all_text(presentation_file)

            return presentation_name, presentation_path, slide_number, hidden_slides, presentation_created, presentation_modified, all_text

        except PackageNotFoundError:
            print(f"Presentation skipped because of password protection: {
                  presentation_selection}")
            return None, None, None, None, None, None, None

    else:
        print("Could not open file")
        return None, None, None, None, None, None, None


def get_time(selection: int, file: str) -> str:
    """0 = creation time, 1 = 'last modified' time"""
    if selection == 1:
        time = os.path.getmtime(file)
    elif selection == 0:
        time = os.path.getctime(file)
    dt = datetime.fromtimestamp(time)
    formatted_time = dt.strftime("%d.%m.%Y %H:%M:%S")

    return formatted_time


def get_hidden_slides(presentation) -> int:
    hidden_slide_counter: int = 0
    for slide in presentation.slides:
        if slide.element.get("show"):
            hidden_slide_counter += 1

    return hidden_slide_counter


def get_all_text(presentation) -> list[str]:
    '''Extract al (text) content & remove all clutter'''
    text_list: list = []
    # Starting the slide counter at 0 because that's what the pptx index function needs
    slide_number: int = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for text in paragraph.runs:
                    # Split to individual words
                    split_text = text.text.split(" ")
                    for split in split_text:
                        # Filtering out all non-alphabetical or numerical characters
                        word: str = ''.join(
                            char for char in split if char.isalnum())
                        # Only words that are not empty strings, have at least 3 characters and are not fully numerical
                        if word != "" and len(word) > 2 and any(char.isalpha() for char in word):
                            text_position: dict = {}
                            text_position["word"] = word
                            text_position["slide"] = slide_number
                            text_list.append(text_position)
        slide_number += 1
    return text_list
