import os
from pptx import Presentation


current_dir: str = os.path.dirname(__file__)
template_path: str = os.path.join(current_dir, "template.pptx")
FILENAME: str = "search_results.pptx"


def create_presentation(found_list) -> None:

    # If a presentation with results already exists - remove it
    if os.path.isfile(FILENAME):
        os.remove(FILENAME)

    prs = Presentation(template_path)

    for found_word in found_list:

        word: str = found_word[0]
        results: list[str] = found_word[1:]

        """Create a slide"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = word  # type: ignore

        """Iterate through the list of found results and fill placeholders accordingly"""
        # Highest placeholder ID: 18; Placeholder Name: Text Placeholder 17
        # Starting at 3 because that's the id of the first placeholder
        placeholder_counter: int = 3
        for result in results:
            # Adjust path from database to windows syntax
            path: str = result[0].replace("/", "\\")

            slide_number: int = int(result[2]) + 1
            result_text: str = f"Found in '{
                result[1]}' on slide {slide_number}"
            result_link: str = f"{path}\\{result[1]}#{slide_number}"

            for placeholder in slide.placeholders:
                if placeholder.shape_id == placeholder_counter:
                    text_frame = placeholder.text_frame  # type: ignore
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = result_text
                    run.hyperlink.address = result_link
            placeholder_counter += 1

    prs.save(FILENAME)
