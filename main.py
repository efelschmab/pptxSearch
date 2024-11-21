from ast import List
import tkinter
import tkinter.filedialog
from pptx import Presentation
import pptx
from customtkinter import filedialog
import sqlite3
import os
from datetime import datetime

# Get current time for timestamps
now = datetime.now()
current_time = now.strftime('%d.%m.%Y %H:%M:%S')

def main() -> None:
    directory_path: str = filedialog.askdirectory()
    found_presentations = find_pptx_files(directory_path)
    for presentation in found_presentations:
        # All seven return values from the get_presentation function put into variables
        presentation_name, presentation_path, slide_number, hidden_slides, presentation_created, presentation_modified, all_text = get_presentation(presentation)
        # Insert the extracted words into the database
        write_to_database(word_list=all_text, presentation_name=presentation_name, presentation_path=presentation_path)
        # Insert the presentation metadata into the database
        save_pptx_to_database(name=presentation_name, path=presentation_path, created=presentation_created, modified=presentation_modified, slides=slide_number, hidden=hidden_slides)
    
    search_query: str = input("Enter your search query: ")
    query_database(search_query)


def find_pptx_files(directory) -> list[str]:
    # This might be a good place to set up some filters, for damaged presentations for example ("._" in front of the filename)
    # Also the check for already existing presentations in the database (in combination with "last modified") should happen here. Skip those existing and not modified to save some time.
    # Bug: not all presentations get recognized (16 out of 30 testpresentations made it into the presentation table)
    file_paths = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith(".pptx"):
                file_path = os.path.join(root, file)
                file_paths.append(file_path)
    return file_paths


def get_presentation(presentation_selection):
    """This is where the file gets opened and every necessary information extracted"""
    
    if presentation_selection:
        # To avoid issues with opening / closing files the presentation is only really opened here
        with open(presentation_selection, 'rb') as presentation_file:
            presentation_file = Presentation(presentation_selection)
            
            # Gathering all metadata
            presentation_name: str = os.path.split(presentation_selection)[1]
            presentation_path: str = os.path.split(presentation_selection)[0]

            slide_number: int = len(presentation_file.slides)
            hidden_slides: int = get_hidden_slides(presentation_file)

            presentation_created: str = get_creation_modified(presentation_file)[0]
            presentation_modified: str = get_creation_modified(presentation_file)[1]

            all_text: list = get_all_text(presentation_file)

        return presentation_name, presentation_path, slide_number, hidden_slides, presentation_created, presentation_modified, all_text
    
    else:
        print("Could not open file")
        return None, None, None, None, None, None, None


def get_hidden_slides(presentation) -> int:
    hidden_slide_counter: int = 0
    for slide in presentation.slides:
        if slide.element.get("show"):
            hidden_slide_counter += 1
    
    return hidden_slide_counter


def get_creation_modified(presentation) -> tuple[str, str]:
   
    created_time: str = "00.00.0000 00:00:00"
    modified_time: str = "00.00.0000 00:00:00"

    if presentation.core_properties.created:
        created_time: str = presentation.core_properties.created.strftime(
            "%d.%m.%Y %H:%M:%S")

    if presentation.core_properties.modified:
        modified_time: str = presentation.core_properties.modified.strftime(
            "%d.%m.%Y %H:%M:%S")

    return created_time, modified_time

def get_all_text(presentation) -> list[str]:
    '''Strip the (text) content of the presentation of all clutter (e.g. empty spaces) and prepare it to be indexed'''
    text_list: list = []   
    # starting the slide counter at 0 because that's what the pptx index function needs
    slide_number: int = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for text in paragraph.runs:
                    # split to individual words
                    split_text = text.text.split(" ")
                    for split in split_text:
                        # filtering out all non-alphabetical or numerical characters
                        word: str = ''.join(char for char in split if char.isalnum())
                        if word != "" and len(split) > 1:
                            text_position: dict = {}
                            text_position["word"] = word
                            text_position["slide"] = slide_number
                            text_list.append(text_position)
        slide_number += 1
    return text_list

def write_to_database(word_list, presentation_name, presentation_path) -> None:
    # Database connection
    database: str = "pptx_search.db"
    conn = sqlite3.connect(database)
    cursor = conn.cursor()

    for word in word_list:
        current_word: str = word['word']
        current_slide_number: int = word["slide"]

        # Check if the table for the word already exists. If so then insert a new row, if not then create it first and then insert.
        cursor.execute(f"SELECT COUNT(*) FROM sqlite_master WHERE type='table' AND name='{current_word}'")
        table_result = cursor.fetchone()
        table_exists = table_result[0] > 0

        sql_insert: str = f"INSERT INTO '{current_word}' (presentation_path, presentation_name, slide_number) VALUES ('{presentation_path}', '{presentation_name}', '{current_slide_number}')"
        
        if table_exists:
            cursor.execute(sql_insert)
        else:
            sql_build_table: str = f"CREATE TABLE IF NOT EXISTS '{current_word}' (presentation_path TEXT, presentation_name TEXT, slide_number INTEGER)"
            cursor.execute(sql_build_table)
            cursor.execute(sql_insert)
    
        # Remove duplicate database entries
        duplicates: str = f"DELETE FROM '{current_word}' WHERE ROWID NOT IN (SELECT MIN(ROWID) FROM '{current_word}' GROUP BY presentation_name, slide_number)"
        cursor.execute(duplicates)
    
    conn.commit()
    conn.close()

def save_pptx_to_database(name, path, created, modified, slides, hidden) -> None:
    # Database connection
    database: str = "pptx_search.db"
    conn = sqlite3.connect(database)
    cursor = conn.cursor()

    sql_build_pptx_table: str = "CREATE TABLE IF NOT EXISTS 'all_presentations' (presentation_name TEXT, presentation_path TEXT, created TEXT, last_modified TEXT, last_indexed TEXT, total_slides INTEGER, hidden_slides INTEGER)"
    cursor.execute(sql_build_pptx_table)

    sql_insert_presentation: str = f"INSERT INTO 'all_presentations' (presentation_name, presentation_path, created, last_modified, last_indexed, total_slides, hidden_slides) VALUES ('{name}', '{path}', '{created}', '{modified}', '{current_time}', '{slides}', '{hidden}')"
    cursor.execute(sql_insert_presentation)
    
    conn.commit()
    conn.close()

def query_database(query) -> None:
    # Database connection
    database: str = "pptx_search.db"
    conn = sqlite3.connect(database)
    cursor = conn.cursor()

    # Get the table (word) names
    sql_search_tables = f"SELECT name FROM sqlite_master WHERE type='table' AND name LIKE '%{query}%'"
    cursor.execute(sql_search_tables)

    table_names = cursor.fetchall()

    for table_name in table_names:
        sql_search_data = f"SELECT presentation_name, slide_number FROM {table_name[0]}"
        cursor.execute(sql_search_data)

        results = cursor.fetchall()
        print(f"{table_name[0]}")
        for row in results:
            print(row)

    conn.close()

# TO-DO: Compare all presentations from the presentations table to the ones on the systems -> update if necessary

# TO-DO: Add word filters and exceptions

# TO-DO: GUI

if __name__ == '__main__':
    main()